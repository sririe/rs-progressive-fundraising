from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUTPUT = Path(
    "/Users/spencer/projects-work/rs-progressive-fundraising/projects/gift-cards/docs/plans/inbox/2026-03-19-discovery-prework-sheet.pptx"
)


RED = RGBColor(220, 38, 38)
DARK = RGBColor(34, 34, 34)
MUTED = RGBColor(95, 99, 104)
LIGHT = RGBColor(245, 245, 245)


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.2), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.08), Inches(11), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED


def add_bullets(slide, items, left=0.95, top=1.7, width=10.2, height=4.8, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)
        p.bullet = True


def add_two_column_bullets(slide, left_items, right_items, left_title, right_title):
    add_title(slide, left_title + " / " + right_title)

    left_box = slide.shapes.add_shape(
        1, Inches(0.7), Inches(1.6), Inches(5.45), Inches(5.3)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = LIGHT
    left_box.line.color.rgb = LIGHT

    right_box = slide.shapes.add_shape(
        1, Inches(6.1), Inches(1.6), Inches(5.45), Inches(5.3)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = LIGHT
    right_box.line.color.rgb = LIGHT

    left_header = slide.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(4.8), Inches(0.4))
    p = left_header.text_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED

    right_header = slide.shapes.add_textbox(Inches(6.4), Inches(1.85), Inches(4.8), Inches(0.4))
    p = right_header.text_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED

    add_bullets(slide, left_items, left=1.0, top=2.25, width=4.8, height=4.2, font_size=16)
    add_bullets(slide, right_items, left=6.4, top=2.25, width=4.8, height=4.2, font_size=16)


def add_footer(slide, text="Progressive Fundraising"):
    line = slide.shapes.add_shape(1, Inches(0.7), Inches(6.8), Inches(11.1), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(225, 225, 225)
    line.line.color.rgb = RGBColor(225, 225, 225)

    footer_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.87), Inches(4), Inches(0.25))
    p = footer_box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def base_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_footer(slide)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(12.33)
    prs.slide_height = Inches(7.0)

    slide = base_slide(prs)
    add_title(slide, "Digital Vendor and Workflow Prework Sheet", "Importable prework deck for the Progressive discovery session")
    hero = slide.shapes.add_textbox(Inches(0.75), Inches(1.95), Inches(11), Inches(2.2))
    tf = hero.text_frame
    p = tf.paragraphs[0]
    p.text = "Purpose"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED
    p = tf.add_paragraph()
    p.text = (
        "Use this deck to align on what we need to understand before and during the on-site session. "
        "The focus is digital fulfillment first, with stickering and invoicing captured as secondary inputs."
    )
    p.font.size = Pt(24)
    p.font.color.rgb = DARK
    p.space_before = Pt(6)
    p.space_after = Pt(14)
    p.alignment = PP_ALIGN.LEFT

    slide = base_slide(prs)
    add_title(slide, "What The Session Must Uncover")
    add_bullets(
        slide,
        [
            "How digital card data is acquired for each vendor",
            "What happens when inventory is received, stored, and prepared for fulfillment",
            "What Lloyd and Mario do step by step for each major vendor flow",
            "Which tools, scripts, spreadsheets, portals, or manual workarounds are involved",
            "Where the workflow is slow, brittle, error-prone, or dependent on operator judgment",
        ],
    )

    slide = base_slide(prs)
    add_two_column_bullets(
        slide,
        [
            "Where card data exists at each stage",
            "Who can access it",
            "How it is transferred or stored today",
            "Whether activation belongs to the digital workflow, the physical workflow, or both",
            "What protections exist today around files, URLs, and customer delivery",
        ],
        [
            "What Mario can do independently today",
            "What still lives mostly in Lloyd's head",
            "What a non-technical operator would still struggle to take over",
            "Which parts could be documented immediately in a runbook",
        ],
        "Sensitive Data And Security",
        "Transferability And Handoffs",
    )

    slide = base_slide(prs)
    add_two_column_bullets(
        slide,
        [
            "Whether end-recipient distribution is a near-term need or still more of an attractive idea",
            "Which delivery expectations are real versus sales-claim assumptions",
            "Whether certain customers or vendors are handled differently today",
        ],
        [
            "High-level understanding of physical-card stickering",
            "High-level understanding of invoicing and QuickBooks friction",
            "Enough context to prioritize whether either one should come after digital fulfillment work",
        ],
        "Delivery Expectations",
        "Adjacent Operational Inputs",
    )

    slide = base_slide(prs)
    add_title(slide, "Prework To Have Ready")
    add_bullets(
        slide,
        [
            "List of active digital vendors",
            "One or two representative workflows ready to walk through live",
            "Sample inputs and outputs where possible",
            "Any current spreadsheets, checklists, scripts, portals, or templates involved in fulfillment",
            "Quick explanation of where digital inventory comes from today",
            "Any constraints around screenshots, recording, or file sharing",
        ],
        font_size=17,
    )

    slide = base_slide(prs)
    add_title(slide, "Framing And Guardrails")
    add_bullets(
        slide,
        [
            "Do not over-architect before seeing the workflow",
            "Treat the session like a walkthrough, not a rigid workshop",
            "Separate customer-facing digital card types from upstream sourcing workflows",
            "Assume PDF/card-generation scripts may be brittle and may not be worth preserving",
            "Record the session if Progressive is comfortable so details can be revisited later",
        ],
        font_size=17,
    )

    slide = base_slide(prs)
    add_title(slide, "Working Questions To Capture Live")
    add_bullets(
        slide,
        [
            "Which digital vendors should be used as representative walkthroughs?",
            "Can Progressive share copies or screenshots of scripts, spreadsheets, sample files, and outputs?",
            "Is recording the full walkthrough allowed?",
            "Does activation belong inside the digital workflow being scoped or mostly inside physical-card operations?",
            "How much of the larger digital roadmap is already being shaped by Lloyd's vendor integration conversations?",
        ],
        font_size=17,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
